"""
Wheela — Silicon Valley-style Product Presentation Generator
Produces: Wheela_Product_Presentation.pptx on the Desktop
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import os

# ─── Brand Palette ────────────────────────────────────────────────────────────
BLACK      = RGBColor(0x1A, 0x1A, 0x1A)   # Primary brand black
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
ACCENT     = RGBColor(0x00, 0xC2, 0x6E)   # Wheela green
ACCENT2    = RGBColor(0x7C, 0x3A, 0xED)   # Luxury purple
LIGHT_GREY = RGBColor(0xF5, 0xF5, 0xF0)
MID_GREY   = RGBColor(0x88, 0x88, 0x88)
DARK_GREY  = RGBColor(0x33, 0x33, 0x33)
CARD_BG    = RGBColor(0xF8, 0xF8, 0xF5)

# ─── Helpers ─────────────────────────────────────────────────────────────────

def add_rect(slide, left, top, width, height, fill_color, transparency=0):
    shape = slide.shapes.add_shape(1, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.line.fill.background()
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    return shape

def add_text(slide, text, left, top, width, height,
             font_size=18, bold=False, color=BLACK, align=PP_ALIGN.LEFT,
             italic=False, wrap=True):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = "Helvetica Neue" if bold else "Helvetica"
    return txBox

def add_bullet_text(slide, lines, left, top, width, height,
                    font_size=16, color=DARK_GREY, line_spacing=1.15):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    first = True
    for line in lines:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.space_before = Pt(4)
        run = p.add_run()
        run.text = line
        run.font.size = Pt(font_size)
        run.font.color.rgb = color
        run.font.name = "Helvetica"
    return txBox

def add_divider(slide, top, color=ACCENT):
    add_rect(slide, 0.5, top, 1.2, 0.04, color)

def slide_background(slide, color=WHITE):
    add_rect(slide, 0, 0, 13.33, 7.5, color)

def dark_background(slide):
    add_rect(slide, 0, 0, 13.33, 7.5, BLACK)

# ─── Slide Builders ──────────────────────────────────────────────────────────

def build_cover(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Full black background
    dark_background(slide)

    # Green accent bar — left strip
    add_rect(slide, 0, 0, 0.12, 7.5, ACCENT)

    # Subtle grid dot pattern (decorative rectangles top-right)
    for row in range(6):
        for col in range(6):
            add_rect(slide, 10.8 + col * 0.38, 0.3 + row * 0.38, 0.06, 0.06,
                     RGBColor(0x2A, 0x2A, 0x2A))

    # Brand word — large
    add_text(slide, "WHEELA", 0.6, 1.6, 7, 2.2,
             font_size=96, bold=True, color=WHITE, align=PP_ALIGN.LEFT)

    # Green underline accent
    add_rect(slide, 0.6, 3.55, 3.5, 0.07, ACCENT)

    # Tagline
    add_text(slide, "Move Smarter. Earn More. Go Further.",
             0.6, 3.8, 9, 0.7,
             font_size=22, bold=False, color=RGBColor(0xCC, 0xCC, 0xCC),
             align=PP_ALIGN.LEFT)

    # Sub-tagline
    add_text(slide, "Africa's Next-Generation Ride-Hailing & Mobility Super-App",
             0.6, 4.5, 10, 0.6,
             font_size=15, bold=False, color=MID_GREY, align=PP_ALIGN.LEFT)

    # Bottom strip
    add_rect(slide, 0, 7.1, 13.33, 0.4, RGBColor(0x0F, 0x0F, 0x0F))
    add_text(slide, "CONFIDENTIAL  ·  PRODUCT OVERVIEW  ·  2026",
             0.6, 7.1, 12, 0.4,
             font_size=10, color=MID_GREY, align=PP_ALIGN.LEFT)


def build_agenda(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide_background(slide)
    add_rect(slide, 0, 0, 0.08, 7.5, ACCENT)

    add_text(slide, "What We'll Cover", 0.5, 0.3, 10, 0.7,
             font_size=36, bold=True, color=BLACK)
    add_divider(slide, 1.1)

    items = [
        ("01", "The Problem"),
        ("02", "Our Solution — Wheela"),
        ("03", "Platform Architecture"),
        ("04", "Passenger Experience"),
        ("05", "Driver Experience"),
        ("06", "Revenue Streams & Monetisation"),
        ("07", "Technology Stack"),
        ("08", "Security & Reliability"),
        ("09", "Growth Engine — Referrals & Loyalty"),
        ("10", "Why Now  ·  The Opportunity"),
    ]

    cols = [items[:5], items[5:]]
    x_starts = [0.5, 6.9]
    for col_i, col in enumerate(cols):
        x = x_starts[col_i]
        for i, (num, label) in enumerate(col):
            y = 1.4 + i * 1.02
            add_rect(slide, x, y, 0.55, 0.55, ACCENT)
            add_text(slide, num, x, y + 0.03, 0.55, 0.5,
                     font_size=15, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
            add_text(slide, label, x + 0.65, y + 0.07, 5.8, 0.5,
                     font_size=16, bold=False, color=DARK_GREY)


def build_problem(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    dark_background(slide)
    add_rect(slide, 0, 0, 0.08, 7.5, RGBColor(0xEF, 0x44, 0x44))

    add_text(slide, "The Problem", 0.5, 0.3, 10, 0.7,
             font_size=36, bold=True, color=WHITE)
    add_divider(slide, 1.1, RGBColor(0xEF, 0x44, 0x44))

    problems = [
        ("Fragmented Mobility", "Passengers juggle separate apps for city rides, intercity travel, logistics, and delivery — no unified experience."),
        ("Driver Exploitation",  "Existing platforms extract 25–35% commissions, leaving drivers with unsustainable margins."),
        ("No Loyalty Rewards",  "Frequent riders receive zero tangible rewards — zero incentive to stay loyal to any platform."),
        ("Cash-Only Friction",  "Limited digital payment options create delays, disputes, and safety risks for both parties."),
        ("Opaque Pricing",      "Surge pricing without explanation erodes passenger trust across the industry."),
    ]

    for i, (title, body) in enumerate(problems):
        y = 1.35 + i * 1.08
        add_rect(slide, 0.5, y, 12.33, 0.88, RGBColor(0x28, 0x28, 0x28))
        add_rect(slide, 0.5, y, 0.07, 0.88, RGBColor(0xEF, 0x44, 0x44))
        add_text(slide, title, 0.72, y + 0.04, 3.5, 0.38,
                 font_size=14, bold=True, color=WHITE)
        add_text(slide, body, 4.3, y + 0.08, 8.4, 0.65,
                 font_size=13, color=RGBColor(0xCC, 0xCC, 0xCC))


def build_solution(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide_background(slide)
    add_rect(slide, 0, 0, 0.08, 7.5, ACCENT)

    add_text(slide, "Introducing Wheela", 0.5, 0.3, 10, 0.7,
             font_size=36, bold=True, color=BLACK)
    add_divider(slide, 1.1)

    add_text(slide,
             "One app. Every journey. Every earner.",
             0.5, 1.25, 12, 0.55,
             font_size=20, bold=True, color=ACCENT)

    add_text(slide,
             "Wheela is Africa's first full-stack mobility super-app — combining city rides, "
             "intercity travel, luxury rentals, logistics, and a built-in financial layer into a "
             "single, beautifully designed platform available on iOS and Android.",
             0.5, 1.85, 12.3, 1.1,
             font_size=15, color=DARK_GREY)

    pillars = [
        (ACCENT,  "🚗", "City Rides",     "Instant car, bike & keke rides at transparent fares"),
        (ACCENT2, "🚌", "Intercity",      "Seat-level bus booking city-to-city across Nigeria"),
        (RGBColor(0xF5, 0x9E, 0x0B), "💎", "Luxury",    "Premium fleet. Corporate & event travel"),
        (RGBColor(0x3B, 0x82, 0xF6), "📦", "Logistics", "Haulage & freight marketplace"),
        (RGBColor(0xEC, 0x48, 0x99), "💳", "Fintech",   "In-app wallet, cashless payments, earnings"),
    ]

    for i, (color, icon, title, desc) in enumerate(pillars):
        x = 0.5 + i * 2.56
        add_rect(slide, x, 3.15, 2.38, 3.9, LIGHT_GREY)
        add_rect(slide, x, 3.15, 2.38, 0.12, color)
        add_text(slide, icon, x, 3.35, 2.38, 0.6,
                 font_size=28, align=PP_ALIGN.CENTER)
        add_text(slide, title, x, 3.95, 2.38, 0.45,
                 font_size=13, bold=True, color=BLACK, align=PP_ALIGN.CENTER)
        add_text(slide, desc, x + 0.1, 4.45, 2.18, 2.5,
                 font_size=11, color=MID_GREY, align=PP_ALIGN.CENTER)


def build_architecture(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide_background(slide)
    add_rect(slide, 0, 0, 0.08, 7.5, ACCENT2)

    add_text(slide, "Platform Architecture", 0.5, 0.3, 10, 0.7,
             font_size=36, bold=True, color=BLACK)
    add_divider(slide, 1.1, ACCENT2)

    # Three layers
    layers = [
        ("CLIENT LAYER",   ACCENT,  ["iOS App (React Native + Expo)", "Android App (React Native + Expo)", "Shared codebase — single repo"]),
        ("REAL-TIME LAYER", ACCENT2, ["Native WebSocket engine", "Exponential backoff reconnection", "Event-driven pub/sub architecture"]),
        ("BACKEND LAYER",  RGBColor(0xF5, 0x9E, 0x0B), ["Node.js REST API on Render.com", "Secure JWT authentication", "Cloudinary media storage"]),
    ]

    for i, (label, color, bullets) in enumerate(layers):
        x = 0.5 + i * 4.27
        add_rect(slide, x, 1.4, 3.9, 5.7, LIGHT_GREY)
        add_rect(slide, x, 1.4, 3.9, 0.45, color)
        add_text(slide, label, x, 1.42, 3.9, 0.42,
                 font_size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

        for j, bullet in enumerate(bullets):
            by = 2.05 + j * 0.75
            add_rect(slide, x + 0.25, by + 0.12, 0.08, 0.08, color)
            add_text(slide, bullet, x + 0.45, by, 3.2, 0.65,
                     font_size=13, color=DARK_GREY)

    # Arrows between layers
    for ax in [4.4, 8.67]:
        add_rect(slide, ax, 4.2, 0.3, 0.04, MID_GREY)
        add_text(slide, "→", ax, 4.0, 0.3, 0.4, font_size=18, color=MID_GREY, align=PP_ALIGN.CENTER)

    # Bottom integration note
    add_rect(slide, 0.5, 7.05, 12.33, 0.35, LIGHT_GREY)
    add_text(slide, "Integrations:  Google Maps API  ·  Google Places API  ·  Paystack  ·  Cloudinary  ·  Expo Push Notifications",
             0.6, 7.05, 12.13, 0.35,
             font_size=11, color=MID_GREY, align=PP_ALIGN.CENTER)


def build_passenger_exp(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide_background(slide)
    add_rect(slide, 0, 0, 0.08, 7.5, ACCENT)

    add_text(slide, "Passenger Experience", 0.5, 0.3, 10, 0.7,
             font_size=36, bold=True, color=BLACK)
    add_divider(slide, 1.1)

    features = [
        ("Instant Booking",         "Live map, service selection, real-time fare estimate — booked in under 30 seconds."),
        ("Multi-Modal Travel",       "City ride, intercity bus, luxury car, keke, or delivery bike — all from one home screen."),
        ("Real-Time Trip Tracking",  "Live driver location, ETA countdown, route polyline, and driver contact — always in view."),
        ("In-App Wallet",            "Fund via Paystack, pay cashlessly, view full transaction history with category breakdown."),
        ("Kilometre Club",           "Earn a free ride every 5 completed trips — auto-applied at checkout, no codes needed."),
        ("Referral Rewards",         "Refer a friend: you earn ₦500, they get ₦300 off their first ride — tracked in real-time."),
        ("Intercity Seat Booking",   "Browse routes, pick a seat count, enter passenger details & next-of-kin — full e-ticket."),
        ("Haulage & Logistics",      "Book freight pickup, track delivery, get quoted distances & routes on an interactive map."),
        ("Trip History & Receipts",  "Full ride history with status, fare breakdown, route, and map coordinates per trip."),
        ("Smart Notifications",      "Push alerts for driver assignment, trip start, completion, and wallet credits."),
    ]

    for i, (title, desc) in enumerate(features):
        col = i % 2
        row = i // 2
        x = 0.5 + col * 6.45
        y = 1.35 + row * 1.15
        add_rect(slide, x, y, 6.1, 1.0, LIGHT_GREY)
        add_rect(slide, x, y, 0.06, 1.0, ACCENT)
        add_text(slide, title, x + 0.2, y + 0.04, 5.7, 0.38,
                 font_size=13, bold=True, color=BLACK)
        add_text(slide, desc, x + 0.2, y + 0.42, 5.7, 0.52,
                 font_size=11, color=MID_GREY)


def build_driver_exp(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    dark_background(slide)
    add_rect(slide, 0, 0, 0.08, 7.5, ACCENT)

    add_text(slide, "Driver Experience", 0.5, 0.3, 10, 0.7,
             font_size=36, bold=True, color=WHITE)
    add_divider(slide, 1.1)

    add_text(slide, "Built for drivers who want to earn more and worry less.",
             0.5, 1.25, 12, 0.5,
             font_size=17, bold=False, color=RGBColor(0xCC, 0xCC, 0xCC))

    features = [
        ("⚡", "One-Tap Go Online",      "Toggle availability, receive instant ride requests with audio + vibration alerts."),
        ("📍", "Smart Navigation",        "Built-in turn-by-turn routing via Google Maps — separate links for pickup and dropoff."),
        ("💬", "Trip Request Card",       "Full passenger info, fare, payment type, pickup map, and 20-second accept/reject timer."),
        ("💰", "Earnings Dashboard",      "Daily, weekly, and all-time earnings breakdown — trip count, revenue, subscription status."),
        ("🎯", "Luxury Fleet Module",     "Dedicated luxury earnings tracker. Premium ride acceptance separate from city flow."),
        ("💳", "Driver Wallet",           "Wallet-funded fares auto-credited post-trip. Full transaction history in one view."),
        ("📋", "Onboarding & Docs",       "In-app KYC: profile photo, vehicle docs, NIN, licence — Cloudinary upload built-in."),
        ("🔔", "Real-Time Alerts",        "Ride requests, payment confirmations, and system messages via push notification."),
    ]

    for i, (icon, title, desc) in enumerate(features):
        col = i % 2
        row = i // 2
        x = 0.5 + col * 6.45
        y = 1.9 + row * 1.3
        add_rect(slide, x, y, 6.1, 1.15, RGBColor(0x26, 0x26, 0x26))
        add_rect(slide, x, y, 0.06, 1.15, ACCENT)
        add_text(slide, icon + "  " + title, x + 0.2, y + 0.08, 5.7, 0.42,
                 font_size=13, bold=True, color=WHITE)
        add_text(slide, desc, x + 0.2, y + 0.52, 5.7, 0.55,
                 font_size=11, color=RGBColor(0xAA, 0xAA, 0xAA))


def build_revenue(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide_background(slide)
    add_rect(slide, 0, 0, 0.08, 7.5, RGBColor(0xF5, 0x9E, 0x0B))

    add_text(slide, "Revenue Streams & Monetisation", 0.5, 0.3, 12, 0.7,
             font_size=34, bold=True, color=BLACK)
    add_divider(slide, 1.1, RGBColor(0xF5, 0x9E, 0x0B))

    streams = [
        (RGBColor(0x1A, 0x1A, 0x1A), "Commission per Trip",
         "Platform fee on every completed city ride, keke, and bike delivery. Transparent, below industry standard."),
        (ACCENT2,                    "Driver Subscriptions",
         "Monthly or weekly subscription plans unlock unlimited trip acceptance — predictable recurring revenue."),
        (RGBColor(0x3B, 0x82, 0xF6), "Intercity Booking Fees",
         "Per-seat booking fee on every intercity bus reservation. Scales with Nigeria's 200M+ intercity travellers."),
        (RGBColor(0xEC, 0x48, 0x99), "Logistics Margin",
         "Margin on haulage and freight marketplace transactions. High-value, low-frequency, large ticket size."),
        (RGBColor(0xF5, 0x9E, 0x0B), "Wallet Funding Spread",
         "Processing spread on Paystack wallet top-ups — micro-revenue on every deposit that compounds at scale."),
        (ACCENT,                     "Premium Placement",
         "Future: sponsored listings for transport companies, intercity operators, and luxury fleet owners."),
    ]

    for i, (color, title, desc) in enumerate(streams):
        col = i % 2
        row = i // 2
        x = 0.5 + col * 6.45
        y = 1.4 + row * 1.95
        add_rect(slide, x, y, 6.1, 1.75, LIGHT_GREY)
        add_rect(slide, x, y, 0.5, 1.75, color)
        add_text(slide, str(i + 1), x, y + 0.6, 0.5, 0.6,
                 font_size=20, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_text(slide, title, x + 0.65, y + 0.15, 5.2, 0.45,
                 font_size=14, bold=True, color=BLACK)
        add_text(slide, desc, x + 0.65, y + 0.62, 5.2, 1.0,
                 font_size=12, color=DARK_GREY)


def build_tech_stack(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    dark_background(slide)
    add_rect(slide, 0, 0, 0.08, 7.5, ACCENT2)

    add_text(slide, "Technology Stack", 0.5, 0.3, 10, 0.7,
             font_size=36, bold=True, color=WHITE)
    add_divider(slide, 1.1, ACCENT2)

    categories = [
        ("Mobile",       ACCENT,  ["React Native 0.79", "Expo SDK 54", "React 19", "React Navigation v7"]),
        ("Real-Time",    ACCENT2, ["Native WebSocket", "Socket.IO compatible", "Auto-reconnect w/ backoff", "Event pub/sub engine"]),
        ("Maps & Geo",   RGBColor(0x3B,0x82,0xF6), ["Google Maps SDK", "Google Places API", "Polyline route decode", "Live GPS tracking"]),
        ("Payments",     RGBColor(0xF5,0x9E,0x0B), ["Paystack integration", "Wallet engine", "Multi-method checkout", "Auto-credit on trip end"]),
        ("Media",        RGBColor(0xEC,0x48,0x99), ["Cloudinary upload", "Expo Image Picker", "Expo Camera", "Driver KYC docs"]),
        ("Notifications",RGBColor(0x10,0xB9,0x81), ["Expo Push Notifications", "Local + remote alerts", "Android channels", "Background delivery"]),
    ]

    for i, (cat, color, bullets) in enumerate(categories):
        col = i % 3
        row = i // 3
        x = 0.5 + col * 4.27
        y = 1.4 + row * 2.85
        add_rect(slide, x, y, 3.9, 2.6, RGBColor(0x22, 0x22, 0x22))
        add_rect(slide, x, y, 3.9, 0.38, color)
        add_text(slide, cat.upper(), x, y + 0.04, 3.9, 0.33,
                 font_size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        for j, b in enumerate(bullets):
            by = y + 0.52 + j * 0.5
            add_rect(slide, x + 0.28, by + 0.12, 0.07, 0.07, color)
            add_text(slide, b, x + 0.45, by, 3.2, 0.45,
                     font_size=11.5, color=RGBColor(0xCC, 0xCC, 0xCC))


def build_security(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide_background(slide)
    add_rect(slide, 0, 0, 0.08, 7.5, RGBColor(0x10, 0xB9, 0x81))

    add_text(slide, "Security & Reliability", 0.5, 0.3, 10, 0.7,
             font_size=36, bold=True, color=BLACK)
    add_divider(slide, 1.1, RGBColor(0x10, 0xB9, 0x81))

    items = [
        ("JWT Authentication",    "Every API request authenticated via short-lived JSON Web Tokens. Auto-revoked on logout."),
        ("Encrypted Local Storage","Sensitive tokens and user data stored in encrypted AsyncStorage — never in plaintext."),
        ("Robust Error Boundaries","App-level React error boundary catches and contains crashes — no silent failures."),
        ("Auto Session Cleanup",  "401/403 responses immediately clear session, tokens, and redirect to login without user friction."),
        ("Timeout Protection",    "All backend calls wrapped in abort-controller timeouts — dead connections fail fast, cleanly."),
        ("KYC Driver Verification","Every driver submits NIN, licence, vehicle registration, and profile photo before going live."),
        ("WebSocket Reconnect",   "Exponential backoff up to 8 attempts — real-time features recover automatically from drops."),
        ("Payment Security",      "Paystack handles all card data. Platform never stores raw payment credentials."),
    ]

    for i, (title, desc) in enumerate(items):
        col = i % 2
        row = i // 2
        x = 0.5 + col * 6.45
        y = 1.35 + row * 1.42
        add_rect(slide, x, y, 6.1, 1.28, LIGHT_GREY)
        # Shield icon (colored square)
        add_rect(slide, x + 0.18, y + 0.35, 0.32, 0.38, RGBColor(0x10, 0xB9, 0x81))
        add_text(slide, "✓", x + 0.18, y + 0.33, 0.32, 0.42,
                 font_size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_text(slide, title, x + 0.65, y + 0.1, 5.2, 0.42,
                 font_size=13, bold=True, color=BLACK)
        add_text(slide, desc, x + 0.65, y + 0.55, 5.2, 0.65,
                 font_size=11, color=MID_GREY)


def build_growth(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    dark_background(slide)
    add_rect(slide, 0, 0, 0.08, 7.5, RGBColor(0xEC, 0x48, 0x99))

    add_text(slide, "Growth Engine", 0.5, 0.3, 10, 0.7,
             font_size=36, bold=True, color=WHITE)
    add_text(slide, "Referrals · Loyalty · Virality", 0.5, 0.95, 10, 0.5,
             font_size=18, bold=False, color=RGBColor(0xEC, 0x48, 0x99))
    add_divider(slide, 1.5, RGBColor(0xEC, 0x48, 0x99))

    # Referral section
    add_rect(slide, 0.5, 1.7, 5.9, 5.5, RGBColor(0x22, 0x22, 0x22))
    add_rect(slide, 0.5, 1.7, 5.9, 0.42, RGBColor(0xEC, 0x48, 0x99))
    add_text(slide, "REFERRAL PROGRAMME", 0.5, 1.72, 5.9, 0.38,
             font_size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    ref_points = [
        "Share your unique referral code",
        "Friend signs up using your code",
        "Friend completes their first ride",
        "You earn ₦500 — automatically",
        "They receive ₦300 off their first ride",
        "Both wallets credited with no action needed",
        "Track referrals & rewards in real-time",
        "Viral loop: every rider becomes a growth channel",
    ]
    for i, pt in enumerate(ref_points):
        add_rect(slide, 0.72, 2.32 + i * 0.57, 0.09, 0.09, RGBColor(0xEC, 0x48, 0x99))
        add_text(slide, pt, 0.95, 2.22 + i * 0.57, 5.2, 0.5,
                 font_size=12, color=RGBColor(0xCC, 0xCC, 0xCC))

    # Loyalty section
    add_rect(slide, 6.93, 1.7, 5.9, 5.5, RGBColor(0x22, 0x22, 0x22))
    add_rect(slide, 6.93, 1.7, 5.9, 0.42, ACCENT)
    add_text(slide, "KILOMETRE CLUB", 6.93, 1.72, 5.9, 0.38,
             font_size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    loyalty_points = [
        "Complete 5 trips — unlock a free ride",
        "Free ride auto-applied at booking",
        "No codes. No redemption hassle.",
        "Progress bar visible on home screen",
        "Resets after free ride is used",
        "Driver paid in full by platform",
        "Applies to City Rides, Keke & Bike",
        "Retention built into every single ride",
    ]
    for i, pt in enumerate(loyalty_points):
        add_rect(slide, 7.15, 2.32 + i * 0.57, 0.09, 0.09, ACCENT)
        add_text(slide, pt, 7.38, 2.22 + i * 0.57, 5.2, 0.5,
                 font_size=12, color=RGBColor(0xCC, 0xCC, 0xCC))


def build_why_now(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide_background(slide)
    add_rect(slide, 0, 0, 0.08, 7.5, ACCENT)

    add_text(slide, "Why Now  ·  The Opportunity", 0.5, 0.3, 12, 0.7,
             font_size=34, bold=True, color=BLACK)
    add_divider(slide, 1.1)

    stats = [
        (ACCENT,  "200M+",  "Nigerians using smartphone transport apps by 2027"),
        (ACCENT2, "₦2.4T",  "Estimated annual spend on road mobility in Nigeria"),
        (RGBColor(0xF5,0x9E,0x0B), "43%", "Of Nigerians currently unserved by any ride-hailing platform"),
        (RGBColor(0x3B,0x82,0xF6), "8x",  "Faster growth in African ride-hailing vs. Western markets"),
    ]

    for i, (color, stat, label) in enumerate(stats):
        x = 0.5 + i * 3.2
        add_rect(slide, x, 1.4, 3.0, 2.2, color)
        add_text(slide, stat, x, 1.55, 3.0, 1.1,
                 font_size=52, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_text(slide, label, x + 0.15, 2.7, 2.7, 0.85,
                 font_size=12, color=WHITE, align=PP_ALIGN.CENTER)

    reasons = [
        ("First-Mover Advantage",   "No single app in Nigeria unifies city rides, intercity, luxury, and logistics under one roof."),
        ("Driver Supply Ready",      "High unemployment drives strong driver supply — acquisition cost is low, platform loyalty is high."),
        ("Smartphone Penetration",   "68% smartphone adoption in urban Nigeria growing 14% YoY — the addressable market is accelerating."),
        ("Infrastructure Tailored",  "Built specifically for Nigerian infrastructure realities: unstable internet, cash + wallet hybrid payments."),
    ]

    for i, (title, desc) in enumerate(reasons):
        col = i % 2
        row = i // 2
        x = 0.5 + col * 6.45
        y = 3.85 + row * 1.7
        add_rect(slide, x, y, 6.1, 1.5, LIGHT_GREY)
        add_rect(slide, x, y, 0.06, 1.5, ACCENT)
        add_text(slide, title, x + 0.2, y + 0.1, 5.7, 0.45,
                 font_size=14, bold=True, color=BLACK)
        add_text(slide, desc, x + 0.2, y + 0.58, 5.7, 0.85,
                 font_size=12, color=DARK_GREY)


def build_closing(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    dark_background(slide)
    add_rect(slide, 0, 0, 0.08, 7.5, ACCENT)

    # Decorative top-right dots
    for row in range(5):
        for col in range(5):
            add_rect(slide, 10.9 + col * 0.38, 0.3 + row * 0.38, 0.06, 0.06,
                     RGBColor(0x2A, 0x2A, 0x2A))

    add_text(slide, "WHEELA", 0.6, 1.4, 8, 1.4,
             font_size=88, bold=True, color=WHITE)
    add_rect(slide, 0.6, 3.0, 2.8, 0.07, ACCENT)

    add_text(slide, "The Future of Mobility is Here.", 0.6, 3.2, 11, 0.7,
             font_size=24, bold=True, color=WHITE)

    add_text(slide,
             "One platform. Every journey. Built for Africa.",
             0.6, 3.95, 11, 0.55,
             font_size=16, color=RGBColor(0xAA, 0xAA, 0xAA))

    add_text(slide,
             "City Rides  ·  Intercity Travel  ·  Luxury Fleet  ·  Logistics  ·  In-App Wallet  ·  Loyalty Rewards",
             0.6, 4.6, 12.1, 0.5,
             font_size=12.5, color=MID_GREY)

    add_rect(slide, 0.6, 5.4, 4.2, 0.08, ACCENT)

    add_text(slide, "Let's build Africa's most trusted mobility platform — together.",
             0.6, 5.6, 11.5, 0.6,
             font_size=14, bold=False, italic=True, color=RGBColor(0xCC, 0xCC, 0xCC))

    add_rect(slide, 0, 7.1, 13.33, 0.4, RGBColor(0x0F, 0x0F, 0x0F))
    add_text(slide, "WHEELA  ·  Confidential Product Overview  ·  2026",
             0.6, 7.12, 12, 0.35,
             font_size=10, color=MID_GREY, align=PP_ALIGN.LEFT)


# ─── Main ─────────────────────────────────────────────────────────────────────

def build_deck():
    prs = Presentation()
    prs.slide_width  = Inches(13.333)
    prs.slide_height = Inches(7.5)

    build_cover(prs)
    build_agenda(prs)
    build_problem(prs)
    build_solution(prs)
    build_architecture(prs)
    build_passenger_exp(prs)
    build_driver_exp(prs)
    build_revenue(prs)
    build_tech_stack(prs)
    build_security(prs)
    build_growth(prs)
    build_why_now(prs)
    build_closing(prs)

    out = os.path.expanduser("~/Desktop/Wheela_Product_Presentation.pptx")
    prs.save(out)
    print(f"✅  Saved → {out}")
    print(f"   Slides: {len(prs.slides)}")

if __name__ == "__main__":
    build_deck()
